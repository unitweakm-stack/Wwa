import os
import logging
import sys
import json
import asyncio
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes, CallbackQueryHandler
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import google.generativeai as genai
from dotenv import load_dotenv

# .env faylidan yuklash (lokal ishlatish uchun)
load_dotenv()

# Logging sozlamalari
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# --- MAXSUS KALITLAR (500 TA KEY) ---
VALID_KEYS = [
    "71TZ7O8T0I", "G85CPCLFWK", "8E8IN8KS1T", "BFWRXYPPYM", "XA5N9YQWYS", "YGGPCEXQKO", "91363JYJ6N", "XU1EW487UT", "62O6LD6F38", "VT8CDRG5V7",
    "5HAB6KPEH0", "TCK39BARDX", "NQEDGVUTAZ", "Z08QZ13I27", "SQSZ33ARTJ", "5VCX925U91", "JN702GDPMY", "HSJ60DIPHA", "USS34JW1L8", "Y8XJNSB0UA",
    "JEUX3WE92S", "1SDLV3KSA4", "PJVDXAMKJE", "OEDWCVCYNQ", "PKIYNLKZ72", "709SBJ3X9N", "9ZRJL4RFK4", "IC4RZCWJQ1", "DD3962HOCT", "UUPSU70EBI",
    "641GQ6OSR6", "U15L23TSL0", "QPY9QZ2K6V", "W9U4C3YDJY", "VPC7SYE0KQ", "4WKIG76KBG", "GUBUEYRPBU", "BTKYAJX91M", "YZO1XBO4PC", "2Y1NOH64GT",
    "03D1VCLTG1", "8QPOZ0JJQ4", "3U92R0H5UR", "U1FUMQ0NJI", "7W493FGV4Q", "KS54CPI80F", "Y5J9SDZ7Q1", "C207U79VSQ", "U2VJXJ7SUZ", "EJB2N1JW1J",
    "EQXKCZHQOA", "Z3DYZPGQ91", "DVZDPATSWO", "F70CTSN3RA", "6J642M8OUF", "888XX4OAN5", "EB1M4FVK5S", "B4P9ZYRWF9", "BN0SMHKC7W", "UIKCODP16I",
    "5LBQESGLRM", "14FYQ27OTB", "GN71S0O4C8", "44EOGRFM74", "Q3BBPMVR19", "1DQ7IZ79OV", "RXHFDHV2SE", "R8Q0UHZI08", "5VGMGT7ZWB", "BN1WUBUNLP",
    "DYIKQE6X4A", "W73HO0H6WK", "KG9GFY6SRX", "Y5B2XFB3IC", "P3WGVNPINB", "6ILL6IYYUY", "UXGX11KTYP", "I9AM6DL9WW", "HE1ERO7K3L", "55WUYY8I34",
    "Z6DZSJAL1B", "LAMHC49QQ1", "7VQTHN1T6K", "8ZMYPY1MOX", "069IGEUX2P", "DF1A13YIOD", "PU108UT29A", "9IOL1ZRGKP", "QFOD9RV8V2", "TJ70QJY49W",
    "FX9JJJI28U", "CCDEC3Q8EA", "RVTC1V97A6", "7DQ5XJNCVN", "LKSATXSMYG", "MWW7NPIVGR", "UCHQ2E5CHA", "0YQA0D4X8R", "FTBAHAWGQ9", "WGHGFD7IGI",
    "9MBIBWPI38", "GAGW9U22UZ", "VI7EG6WCSJ", "7CZGVJJEWZ", "KAACQF7MPS", "ILDPYUROTS", "3645R6N77B", "27RHT5F7H0", "T7MI4EFGZN", "OLDFIQ2LDA",
    "B448SOA1XG", "N1V6PVEOJO", "8CX8OQ6L91", "Q366DKM2HT", "JYUHQAR6LF", "6P8K8E2KAW", "QDGBXZ9CZ8", "B0PSMDKGUN", "E50YFRSAAK", "7HC4H7LBJR",
    "QRMRVNM2QW", "T5QO6RF8CB", "LYB2UTMIX7", "T1YGQ4B50T", "AAO2EAW6IF", "3K167PO888", "3G7YFQPQ8B", "FKB31OUTH1", "09BUOP523H", "6359FKWOW0",
    "UQINVE89JQ", "E3DMJLGU4C", "ZSFJ08DCUQ", "CCBSAIENMS", "WXE1NS22UT", "K2LC5MM5JQ", "BXTM2R20YD", "T73UO0YJP6", "VJ9J4BDNLX", "8E2S4HX1FH",
    "X56E0RIBR4", "CP654P4MKB", "5JY5PY76ZA", "QP2T4EUIWI", "1ULFHR8HIE", "C3UQHU52D4", "I5N7SQVB7P", "3C7SXF482H", "6YAUTFEV5E", "1IZAU5EUVX",
    "J2P4P3N68Q", "ENN7DJ1R3M", "VVKQFXA7T0", "AOHT5GISBL", "G227D72ETD", "HI96GTSDWE", "1MBPWY1708", "UKHY7U88PG", "I7RNHPMCW2", "A5C072QYUK",
    "RFVOMNSL98", "RZVRTAAU82", "D005XOISZ8", "WX7ZUC3WW0", "2VHIM4H07K", "X6O32HU7U3", "4H2SI8PBI5", "CHOID3CFPF", "R1TGBYK8VT", "6ZO37VI9HA",
    "P3VOHUPIQC", "N6GT3UBOTU", "P8X070UAV4", "2QB1WCMA8J", "OWFX5OYG3Y", "981T7JTUSH", "VSLMCDRIZY", "USKM4ULSPZ", "8NR9NYE1LD", "FOAU0P3YS1",
    "65SJGWLQ0X", "NNMI2YGZ3O", "7N5U2ND44Z", "9MF9R969WH", "NC51I8NV1B", "18OBHJ3MMJ", "0VNCQVCU1N", "2TRV7KYVQX", "IR1T69PCQX", "RXO26JL2IL",
    "ITVLWOD0PF", "X2JNB0UX2T", "058R2QCWO4", "ILTFJGJS08", "O2GNGHZBOL", "KTWMDF4VZJ", "FD5CHDNI1O", "BZZTYQ9FNC", "EUOEI1ASPP", "B5ZRXELA8K",
    "VH97KFGRVP", "H481WXZFRK", "I2SS213XXQ", "WFLEHZDQB5", "H958GL6CJU", "OB14HAE3MQ", "GFOHOVZ4BK", "OTSU1PJUKF", "YYBHV9UPBI", "R8510ZV1SA",
    "ENXDU5T75J", "TF8MRIZ78T", "Z3Y8LTNC3K", "1JMKYR6GQD", "7SL02GN65A", "K89QRQYXZE", "4N4H6RRA9G", "SOSHIMSVRF", "DGJDK7WBM9", "AGLTKNCU2B",
    "AODGMMALSS", "18QEL4M372", "1PYZVX902U", "5UY1VIIVJG", "Y4VPXEUL1L", "1BMDNKL5N4", "TM17CZICKX", "NHPEAXLR6R", "QLVF86FYC8", "22K3W61Y6M",
    "ZLLOA5D839", "NM8M952T0W", "CW526TBRZ9", "NUIXBE8T4A", "X0J5O1X6RQ", "RQVSSDRY97", "GA0SRQ9N56", "KF6NE9EFP9", "7TSFMC3534", "3PLBZ7NVFC",
    "4NDJBM567W", "2X29DY6C5C", "06FVABO9HB", "PMYVWL3KEN", "7OFI297B57", "EMYZ8WG0RM", "726CNWXVA6", "8R9J07TJMO", "ZGCZ462KY0", "MIMPYQHWKL",
    "8K7Q62GH5F", "RYDRTYWKME", "7ETLPE5Q6K", "DQOZIUNN01", "Y0NQTF552H", "756BMLDGEA", "J4R8CC2YR8", "UH84HGUKK5", "E791MKS6H9", "6UT47HMKD0",
    "IWZKFNZ5TH", "U3UGTR8DNA", "N9DI82Q9HO", "X5WZQ3CCRI", "GOMC0TKS6N", "758H3F1ZDW", "UQD8KBFWXK", "RBD71TLRTK", "77UN1JOOAZ", "OEPJIPB75W",
    "JF5RV81QBA", "BHVOCVFZ95", "ZUMFYBI1LU", "4EE14NQ4A5", "JXZEQIIXVU", "L3YL9HEALP", "3B5S9VGNYQ", "VPBZHLD0M9", "II7QUDKL6M", "JBZHAN6CLE",
    "5EZMHGH1TR", "UL09Q55KVN", "2OFZX6NG2Z", "NH7L8ORK12", "77SEA04YCR", "KIED0GWYYQ", "FLXJ1NCCUX", "HB4AEC61W8", "49Z5M6D13I", "XFSIK0IQWI",
    "3I5D0JN1ZD", "YD6R1957JX", "SB0TTC9Z80", "AVVANPSMIZ", "ZMX57KB31U", "S286TEYAGG", "GZ00W19PYW", "NT2CSCA2P2", "BRLKU8Z19Q", "0YU6F2MULT",
    "M4IN4ERBY2", "AHNOIUN6CX", "2G6MZH3PCD", "010ULEMQ7I", "EUTV8VCRSG", "FX8EU451K8", "CUQJH6SJIX", "9KES4UUCBU", "5M0M9BDVHT", "2F7DCB7DFY",
    "OD8NGXRXNV", "21ZC9I7AVX", "PSG6U9N0TR", "JGKKW0JSKR", "DJAQJLGMEK", "CCT8K51N3Z", "J09HP7WY00", "XG196T2AJP", "IV7IDNKJ46", "P1ENK0QDDI",
    "WPBEGEMV3G", "US2H2XX96U", "RYQ98H0EZV", "8RWJZBNBCG", "TX99G0RIQB", "RZKNFNKSZY", "IUPYKQKEY2", "CT70UKNRFX", "H9KBMHNMM5", "0PN8VJTODC",
    "OXOIGYP286", "DWEPB4XKDG", "0WZGHKHT9P", "DBIWHP6BR2", "AZDTWRWZJV", "THF3EWCHK6", "DAJRIN2V08", "5LBDH3081T", "K6TP6CCDZ4", "F8W86X133D",
    "CUWDAJQM34", "LTX2S5QZ4Y", "KBQ6V1Z7IG", "3OQA8I14GH", "BGDJPF348L", "IMM7FEUU1V", "59P01HLN8C", "X1Z1MSRRXV", "210FX74HEN", "QKRP8CT4OV",
    "7WBPLJTUER", "LS0U3VMN61", "5MBQ8HHICY", "FLBKM3U52Z", "928FTPHF26", "3YDCLL5KSS", "TX2H8AVLNT", "IUB8TMR7DI", "49401KNUUB", "V29X3FBV60",
    "GFS7XIAKPW", "8CNV3V1BL7", "7CAI0X5JJJ", "CBQLDJX1XG", "7HZ9ERUDP5", "2JKBICZ6YV", "G9KIWTSOGH", "S3YMKVDTCV", "SP6XIEVI55", "TDOB97RFVO",
    "TAS3R63NAF", "35X8Q46O0A", "FTSUOZD8IX", "MC25BIJ8HY", "JE42PK6H3C", "E2W0WLF4BG", "MMRVCTY3DG", "2EVYGCAU0P", "45HCWR7QDZ", "4QLDDRM4NO",
    "RC5M81B4NM", "JDBJ1S245H", "6Y28E0J02V", "RC1U3LF0CT", "B97H6SGLZN", "PDRZXNA4LS", "B212BBUK20", "HZE6VX026G", "3NTDF0YX10", "PYQSP5HMO2",
    "DEPI4WBE5E", "1PASGM85FY", "YGCLKJ3S1V", "H4LT8H66ND", "N7G034QVFD", "YF06I7FLVB", "IC5S75H12U", "2BY1F76S2B", "CDD0K0NW3D", "UPLFOKWTN8",
    "Q2TT9R6D7W", "03SRX0EJE5", "66ICMMDLY2", "IIDHS5RAIE", "42A4QA8FQM", "3721G7MB0Z", "Y2AZ2K203A", "7YO29REMCA", "Y2HUJ1U4NS", "XZE5X4GAAW",
    "K989RS9CFD", "OXSRL998KA", "PT44G23UN2", "KT8R2JIY6F", "2LYG03WDYJ", "3FVGEHVFGJ", "HIQINRUGUZ", "OILVYY0GSW", "ABRWCT9H39", "XTIBWDOMWK",
    "K5MFWF3GWM", "E0X6C0ZYBN", "YY4ZTSLQNA", "C3V17L1SX0", "OGOSPQHFK8", "N6D5HBOLPF", "A4KCGJPCRX", "ROFP1856D8", "X3LXJPBRKA", "OPBLGL8N3W",
    "I6URZPCJK9", "P3UDMXX6KS", "A7ZWFP1MJB", "MA1MJGQW0O", "3P2PK3AL88", "JORMDKZ3E8", "XZ7ZJEPE8P", "I5X68E40BR", "FHOFWPP1VH", "OXN47URM69",
    "OOMAQG1907", "UJGIQRU2KA", "0PK7443Q37", "IAUSGPQ6M0", "5S9V1EB7LY", "HHMRNLHPF8", "68CUUSVQ80", "IM0FE339KN", "38UOMWCIXV", "OKE6ZC4LLU",
    "4FXGSDPFU0", "T3R0HOYU7D", "QWS9LPV1YX", "XNTDAR3SYC", "Z2EOQ98VPC", "IX18R1YSRS", "WI9G627D7U", "HU7LAKWC4J", "ZNRVV1LESN", "PAO93RY6U5",
    "HP2MKR26XM", "JUW3ZMXMMZ", "EA7R9SWQYV", "09L1X1889U", "Y0Q7K6K3K7", "X91OK03C48", "FCQ5R7IRFE", "VKEUWF0UVK", "CX3TVCHCP7", "R47TLD6GYW",
    "S896TK9QS1", "WT6WWT2EE4", "4FUHSW35CY", "SJ5ZRXYUJE", "2YA5A7FRL7", "LP3M7DX129", "NWY37N31IX", "PWA2R6XM1W", "3J7BJVELU5", "P2GY6HVBWR",
    "B8Z0ZF9N3A", "QJOS5OU55Z", "EFU937MNAL", "WY8M5YGY03", "5LNJOR3UW3", "3YQXIHNGUJ", "P7DEDO1IQ6", "JFF6UL4GU7", "HRAJPB76N9", "WELLRLWFS4",
    "73W78WHRTA", "DFZB9OAW8G", "EZ5LJWQ2UQ", "LMB6G5TIIU", "T46U8RJ63L", "SHZ6890TTY", "P55KVBDEFM", "72U8FRYKBS", "OVR5MA6XV9", "607DWS7Z7Q"
]

# --- SLAYD GENERATOR QISMI ---
class SlideGenerator:
    def __init__(self, api_key=None):
        self.themes = {
            "akademik": {"bg_color": RGBColor(255, 255, 255), "title_color": RGBColor(0, 51, 102), "text_color": RGBColor(0, 0, 0), "font_name": "Arial"},
            "modern": {"bg_color": RGBColor(240, 240, 240), "title_color": RGBColor(204, 0, 0), "text_color": RGBColor(50, 50, 50), "font_name": "Calibri"},
            "dark": {"bg_color": RGBColor(30, 30, 30), "title_color": RGBColor(255, 255, 255), "text_color": RGBColor(200, 200, 200), "font_name": "Verdana"}
        }
        if api_key and api_key != "YOUR_GEMINI_API_KEY":
            try:
                genai.configure(api_key=api_key)
                self.model = genai.GenerativeModel('gemini-1.5-flash')
            except Exception as e:
                logger.error(f"Gemini AI xatosi: {e}")
                self.model = None
        else:
            self.model = None

    def generate_content_with_ai(self, topic):
        if not self.model:
            return [
                {"title": "Kirish", "content": f"{topic} haqida umumiy ma'lumot."},
                {"title": "Asosiy qism", "content": f"{topic}ning ahamiyati va bugungi kundagi o'rni."},
                {"title": "Xulosa", "content": "Mavzu bo'yicha yakuniy xulosalar."}
            ]
        
        prompt = f"Mavzu: {topic}. Institut talabalari uchun 5-7 ta slayddan iborat mazmun yarat. Javobni faqat JSON formatida qaytar: [{{'title': '...', 'content': '...'}}]. Til: O'zbekcha."
        try:
            response = self.model.generate_content(prompt)
            text = response.text
            if "```json" in text: text = text.split("```json")[1].split("```")[0]
            elif "```" in text: text = text.split("```")[1].split("```")[0]
            return json.loads(text.strip())
        except Exception as e:
            logger.error(f"AI Content Error: {e}")
            return [{"title": "Kirish", "content": f"{topic} haqida ma'lumot."}, {"title": "Xulosa", "content": "Tugadi."}]

    def create_presentation(self, topic, content_list, theme_name, output_pptx):
        prs = Presentation()
        theme = self.themes.get(theme_name, self.themes["akademik"])
        
        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = topic
        slide.placeholders[1].text = "Institut talabalari uchun\nGemini AI orqali yaratildi"

        # Content Slides
        for item in content_list:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = item.get('title', 'Mavzu')
            tf = slide.placeholders[1].text_frame
            tf.text = item.get('content', '')
            for p in tf.paragraphs:
                p.font.size = Pt(18)
                p.font.name = theme['font_name']
                p.font.color.rgb = theme['text_color']
        
        prs.save(output_pptx)
        return output_pptx

# --- TELEGRAM BOT QISMI ---
BOT_TOKEN = os.getenv("BOT_TOKEN", "YOUR_TELEGRAM_BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "YOUR_GEMINI_API_KEY")
user_data = {}

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    user_data[user_id] = {"authenticated": False, "step": "auth", "slides_count": 0}
    await update.message.reply_text("🎓 Slayd botga xush kelibsiz!\nBotni ishlatish uchun maxsus kalitni (KEY) kiriting:")

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    text = update.message.text
    if user_id not in user_data: user_data[user_id] = {"authenticated": False, "step": "auth", "slides_count": 0}

    if not user_data[user_id]["authenticated"]:
        if text in VALID_KEYS:
            user_data[user_id].update({"authenticated": True, "step": "topic"})
            await update.message.reply_text("✅ Tasdiqlandi! Endi slayd mavzusini yuboring:")
        else:
            await update.message.reply_text("❌ Noto'g'ri kalit! Iltimos, to'g'ri kalitni kiriting.")
        return

    if user_data[user_id]["step"] == "topic":
        user_data[user_id].update({"topic": text, "step": "theme"})
        keyboard = [[InlineKeyboardButton("Akademik", callback_data='akademik')], [InlineKeyboardButton("Modern", callback_data='modern')], [InlineKeyboardButton("Dark", callback_data='dark')]]
        await update.message.reply_text(f"Mavzu: {text}\nDizaynni tanlang:", reply_markup=InlineKeyboardMarkup(keyboard))

async def button_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    user_id = query.from_user.id
    theme = query.data
    await query.answer()
    await query.edit_message_text(text="Gemini AI slayd tayyorlamoqda, kuting...")

    topic = user_data[user_id].get("topic", "Mavzu")
    generator = SlideGenerator(api_key=GEMINI_API_KEY)
    content = generator.generate_content_with_ai(topic)
    pptx_file = f"slide_{user_id}.pptx"
    generator.create_presentation(topic, content, theme, pptx_file)

    try:
        with open(pptx_file, 'rb') as doc:
            await context.bot.send_document(chat_id=user_id, document=doc, filename=f"{topic}.pptx", caption=f"✅ {topic} tayyor!")
        user_data[user_id]["slides_count"] += 1
        if user_data[user_id]["slides_count"] >= 2:
            user_data[user_id].update({"authenticated": False, "slides_count": 0, "step": "auth"})
            await context.bot.send_message(chat_id=user_id, text="⚠️ Limit tugadi. Yangi kalit kiriting:")
        else:
            user_data[user_id]["step"] = "topic"
            await context.bot.send_message(chat_id=user_id, text="Yana 1 ta slayd yaratishingiz mumkin. Mavzu yuboring:")
    except Exception as e:
        await query.message.reply_text(f"Xatolik: {e}")
    finally:
        if os.path.exists(pptx_file): os.remove(pptx_file)

if __name__ == '__main__':
    if not BOT_TOKEN or BOT_TOKEN == "YOUR_TELEGRAM_BOT_TOKEN":
        logger.error("BOT_TOKEN o'rnatilmagan!")
        sys.exit(1)
    app = ApplicationBuilder().token(BOT_TOKEN).build()
    app.add_handler(CommandHandler("start", start))
    app.add_handler(MessageHandler(filters.TEXT & (~filters.COMMAND), handle_message))
    app.add_handler(CallbackQueryHandler(button_callback))
    logger.info("Bot ishga tushdi...")
    app.run_polling()
