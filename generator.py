import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import google.generativeai as genai
import json

class SlideGenerator:
    def __init__(self, api_key=None):
        self.themes = {
            "akademik": {
                "bg_color": RGBColor(255, 255, 255),
                "title_color": RGBColor(0, 51, 102),
                "text_color": RGBColor(0, 0, 0),
                "font_name": "Arial"
            },
            "modern": {
                "bg_color": RGBColor(240, 240, 240),
                "title_color": RGBColor(204, 0, 0),
                "text_color": RGBColor(50, 50, 50),
                "font_name": "Calibri"
            },
            "dark": {
                "bg_color": RGBColor(30, 30, 30),
                "title_color": RGBColor(255, 255, 255),
                "text_color": RGBColor(200, 200, 200),
                "font_name": "Verdana"
            }
        }
        if api_key and api_key != "YOUR_GEMINI_API_KEY":
            try:
                genai.configure(api_key=api_key)
                self.model = genai.GenerativeModel('gemini-1.5-flash')
            except Exception as e:
                print(f"Gemini AI konfiguratsiya xatosi: {e}")
                self.model = None
        else:
            print("GEMINI_API_KEY topilmadi yoki default qiymatda qoldi. Slaydlar Gemini AI-siz yaratiladi.")
            self.model = None

    def generate_content_with_ai(self, topic):
        if not self.model:
            # Agar API key bo'lmasa, namunaviy mazmun qaytarish
            return [
                {"title": "Kirish", "content": f"{topic} haqida umumiy ma'lumot."},
                {"title": "Asosiy qism", "content": f"{topic}ning ahamiyati va bugungi kundagi o'rni."},
                {"title": "Tahlil", "content": "Institut talabalari uchun tahliliy ma'lumotlar."},
                {"title": "Xulosa", "content": "Mavzu bo'yicha yakuniy xulosalar."}
            ]

        prompt = f"""
        Mavzu: {topic}
        Ushbu mavzu bo'yicha institut talabalari uchun slayd mazmunini yaratib ber.
        Mazmun 5-7 ta slayddan iborat bo'lsin.
        Har bir slayd uchun 'title' (sarlavha) va 'content' (matn) bo'lsin.
        Javobni faqat JSON formatida qaytar, mana bu ko'rinishda:
        [
            {{"title": "Sarlavha 1", "content": "Matn 1"}},
            {{"title": "Sarlavha 2", "content": "Matn 2"}}
        ]
        Matnlar o'zbek tilida bo'lsin.
        """
        
        try:
            response = self.model.generate_content(prompt)
            # JSON qismini ajratib olish (ba'zan AI markdown formatida qaytaradi)
            text = response.text
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0]
            elif "```" in text:
                text = text.split("```")[1].split("```")[0]
            
            return json.loads(text.strip())
        except Exception as e:
            print(f"AI Error: {e}")
            return self.generate_content_with_ai(None) # Fallback

    def create_presentation(self, topic, content_list, theme_name="akademik", output_pptx="output.pptx"):
        prs = Presentation()
        theme = self.themes.get(theme_name, self.themes["akademik"])

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic
        subtitle.text = "Institut talabalari uchun maxsus tayyorlandi\nGemini AI orqali yaratildi"

        # Content Slides
        for item in content_list:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = item.get('title', 'Mavzu')
            tf = body_shape.text_frame
            tf.text = item.get('content', '')
            
            # Apply basic styling
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = theme['font_name']
                paragraph.font.color.rgb = theme['text_color']

        prs.save(output_pptx)
        return output_pptx
